import pptx
from pptx.util import Inches
from pptx import Presentation
import glob
import os

#reference url: https://qiita.com/aykbzl/items/09b52fabb3af6b925fb9

ppt = Presentation()
width = ppt.slide_width
height = ppt.slide_height

#使用するスライドの種類
title_slide_layout = ppt.slide_layouts[0] #Title Slideの作成
bullet_slide_layout = ppt.slide_layouts[1] #Title and Contentの作成
blank_slide_layout = ppt.slide_layouts[6] #Blankの作成

#Title Slide
slide = ppt.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

#title_text = input("Title: ")
#subtitle_text = input("Subtitle: ")

title_text = ""
subtitle_text = ""


title.text = title_text
subtitle.text = subtitle_text

#get images
fnms = glob.glob('*.png')
#print(fnms)
tx_left = tx_top = tx_width = tx_height = Inches(1)
i = 1
for fnm in fnms:
    #insert images
    slide_picture = ppt.slides.add_slide(blank_slide_layout)
    pic = slide_picture.shapes.add_picture(fnm, width, height, width, height)
    pic.left = int((width -pic.width)/2)
    pic.top = int((width -pic.height)/2)
    #insert a textbox
    txBox = slide_picture.shapes.add_textbox(tx_left, tx_top/2, tx_width, tx_height)
    tB = txBox.text_frame
    tB.text = "No. " + str(i)
    i += 1
ppt.save('figure_zircon.pptx')
os.system("open figure_zircon.pptx")
